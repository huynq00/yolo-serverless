#!/usr/bin/env python3
"""Generate YOLO Serverless presentation using the official course PPT template style.

Template: Báo Cáo Đồ Án Hệ Tính Toán Phân Bố Nâng Cao.pptx
- Fonts: Merriweather (titles), DM Sans (body)
- Accent: #06B6D4
- Background images from template-assets/
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "report-figures"
ASSETS = ROOT / "template-assets"
TEMPLATE = ROOT / "Báo Cáo Đồ Án Hệ Tính Toán Phân Bố Nâng Cao.pptx"
OUT = ROOT / "Thuyet_Trinh_YOLO_Serverless.pptx"

# Template palette
CYAN = RGBColor(0x06, 0xB6, 0xD4)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
NAVY_BLUE = RGBColor(0x1E, 0x3A, 0x8A)
SLATE_DARK = RGBColor(0x1E, 0x29, 0x3B)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
ROW_LINE = RGBColor(0xE2, 0xE8, 0xF0)

FONT_TITLE = "Merriweather"
FONT_BODY = "DM Sans"
FONT_BODY_MED = "DM Sans Medium"

BG_TITLE = ASSETS / "image1.png"       # dark cover
BG_CONTENT = ASSETS / "image2.png"     # light content
BG_CONCLUSION = ASSETS / "image25.png" # dark conclusion
BULLET_ICON = ASSETS / "image5.png"    # cyan check/dot
CARD_LIGHT = ASSETS / "image6.png"     # light card panel
CARD_DARK = ASSETS / "image7.png"      # dark card panel
CARD_3COL = ASSETS / "image10.png"     # 3-col card
SIDE_PANEL = ASSETS / "image3.png"     # right illustration panel
ICON_QA = ASSETS / "image27.png"


def _set_run_font(run, *, size, bold=False, color=SLATE, font=FONT_BODY):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    # Ensure East Asian fallback
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", font)


def _add_textbox(slide, left, top, width, height, text, *,
                 size=14.25, bold=False, color=SLATE, font=FONT_BODY,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size=size, bold=bold, color=color, font=font)
    return box


def _add_multiline(slide, left, top, width, height, lines, *,
                   size=14.25, bold=False, color=SLATE, font=FONT_BODY,
                   align=PP_ALIGN.LEFT, space_after=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        _set_run_font(run, size=size, bold=bold, color=color, font=font)
    return box


def _add_bg(slide, image_path: Path):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0), Inches(0),
                                 width=Inches(13.333), height=Inches(7.5))


def _title_with_accent(slide, title: str):
    """Cyan vertical bar + Merriweather title — matches template content slides."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(0.62), Inches(0.08), Inches(0.57)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    _add_textbox(
        slide, Inches(0.83), Inches(0.62), Inches(12.0), Inches(0.57),
        title, size=33, bold=True, color=NAVY, font=FONT_TITLE,
    )


def _blank_slide(prs) -> object:
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank


def add_cover(prs):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_TITLE)
    _add_textbox(
        slide, Inches(0.69), Inches(1.85), Inches(11.94), Inches(1.35),
        "Tối Ưu Hóa Cold-Start AI Serverless\nbằng Weight-Sharing",
        size=36, bold=True, color=LIGHT, font=FONT_TITLE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, Inches(0.98), Inches(3.45), Inches(11.38), Inches(0.43),
        "Báo Cáo Đồ Án — Hệ Tính Toán Phân Bố Nâng Cao (NT2204.CH201)",
        size=18, bold=False, color=LIGHT, font=FONT_BODY_MED, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, Inches(0.98), Inches(4.45), Inches(11.38), Inches(0.32),
        "Sinh viên thực hiện: Ngô Quang Huy",
        size=14.25, color=SLATE, font=FONT_BODY, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, Inches(0.98), Inches(4.90), Inches(11.38), Inches(0.32),
        "Giảng viên hướng dẫn: TS. Huỳnh Văn Đặng",
        size=14.25, color=SLATE, font=FONT_BODY, align=PP_ALIGN.CENTER,
    )


def add_section(prs, number: str, title: str, subtitle: str):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _add_textbox(
        slide, Inches(1.5), Inches(2.55), Inches(10.3), Inches(0.85),
        f"{number}. {title}",
        size=42, bold=True, color=NAVY_BLUE, font=FONT_TITLE, align=PP_ALIGN.CENTER,
    )
    # Cyan underline
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.15), Inches(3.55), Inches(1.04), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    _add_textbox(
        slide, Inches(2.5), Inches(4.0), Inches(8.3), Inches(0.8),
        subtitle, size=16, color=SLATE, font=FONT_BODY, align=PP_ALIGN.CENTER,
    )


def add_bullets(prs, title: str, bullets: list[str], intro: str | None = None):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _title_with_accent(slide, title)

    y = 1.45
    if intro:
        _add_textbox(
            slide, Inches(0.62), Inches(y), Inches(12.0), Inches(0.7),
            intro, size=14.25, color=SLATE, font=FONT_BODY,
        )
        y = 2.25

    for bullet in bullets:
        if BULLET_ICON.exists():
            slide.shapes.add_picture(
                str(BULLET_ICON), Inches(0.62), Inches(y + 0.04),
                width=Inches(0.25), height=Inches(0.27),
            )
        _add_textbox(
            slide, Inches(1.04), Inches(y), Inches(11.5), Inches(0.55),
            bullet, size=14.25, color=SLATE, font=FONT_BODY,
        )
        y += 0.58
    return slide


def add_two_cards(prs, title: str,
                  left_title: str, left_body: str, left_bullets: list[str],
                  right_title: str, right_body: str, right_bullets: list[str],
                  left_dark: bool = False):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _title_with_accent(slide, title)

    # Card backgrounds
    left_img = CARD_DARK if left_dark else CARD_LIGHT
    right_img = CARD_DARK if not left_dark else CARD_LIGHT
    # Use light cards for both unless specified; template uses image6/image7
    if CARD_LIGHT.exists():
        slide.shapes.add_picture(str(CARD_LIGHT), Inches(0.62), Inches(1.55),
                                 width=Inches(5.83), height=Inches(5.0))
    if CARD_LIGHT.exists():
        # For contrast, right can reuse same light card or dark variant
        img = CARD_DARK if CARD_DARK.exists() and left_dark is False else CARD_LIGHT
        # Keep both light for readability of long text; use cyan title accents
        slide.shapes.add_picture(str(CARD_LIGHT), Inches(6.88), Inches(1.55),
                                 width=Inches(5.83), height=Inches(5.0))

    # Left column
    _add_textbox(slide, Inches(1.04), Inches(1.85), Inches(5.0), Inches(0.35),
                 left_title, size=16, bold=True, color=CYAN, font=FONT_TITLE)
    _add_textbox(slide, Inches(1.04), Inches(2.30), Inches(5.0), Inches(1.0),
                 left_body, size=13, color=SLATE, font=FONT_BODY)
    y = 3.45
    for b in left_bullets:
        if BULLET_ICON.exists():
            slide.shapes.add_picture(str(BULLET_ICON), Inches(1.04), Inches(y),
                                     width=Inches(0.22), height=Inches(0.24))
        _add_textbox(slide, Inches(1.40), Inches(y - 0.02), Inches(4.6), Inches(0.4),
                     b, size=13, color=SLATE, font=FONT_BODY)
        y += 0.48

    # Right column
    _add_textbox(slide, Inches(7.29), Inches(1.85), Inches(5.0), Inches(0.35),
                 right_title, size=16, bold=True, color=CYAN, font=FONT_TITLE)
    _add_textbox(slide, Inches(7.29), Inches(2.30), Inches(5.0), Inches(1.0),
                 right_body, size=13, color=SLATE, font=FONT_BODY)
    y = 3.45
    for b in right_bullets:
        if BULLET_ICON.exists():
            slide.shapes.add_picture(str(BULLET_ICON), Inches(7.29), Inches(y),
                                     width=Inches(0.22), height=Inches(0.24))
        _add_textbox(slide, Inches(7.65), Inches(y - 0.02), Inches(4.6), Inches(0.4),
                     b, size=13, color=SLATE, font=FONT_BODY)
        y += 0.48
    return slide


def add_three_cards(prs, title: str, cards: list[tuple[str, str]]):
    """cards: [(card_title, card_body), ...] length 3"""
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _title_with_accent(slide, title)

    xs = [0.62, 4.76, 8.89]
    for i, (ct, cb) in enumerate(cards[:3]):
        x = xs[i]
        if CARD_3COL.exists():
            slide.shapes.add_picture(str(CARD_3COL), Inches(x), Inches(1.70),
                                     width=Inches(3.82), height=Inches(4.6))
        _add_textbox(slide, Inches(x + 0.32), Inches(2.55), Inches(3.2), Inches(0.5),
                     ct, size=16, bold=True, color=NAVY, font=FONT_TITLE)
        _add_textbox(slide, Inches(x + 0.32), Inches(3.20), Inches(3.2), Inches(2.5),
                     cb, size=13, color=SLATE, font=FONT_BODY)
    return slide


def add_split_content(prs, title: str, left_title: str, left_intro: str,
                      bullets: list[str], image_path: Path | None = None):
    """Left text + right image panel (template slide 6 style)."""
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _title_with_accent(slide, title)

    _add_textbox(slide, Inches(0.62), Inches(1.50), Inches(5.7), Inches(0.35),
                 left_title, size=16, bold=True, color=SLATE_DARK, font=FONT_TITLE)
    _add_textbox(slide, Inches(0.62), Inches(1.95), Inches(5.5), Inches(0.9),
                 left_intro, size=13.5, color=SLATE, font=FONT_BODY)

    y = 3.10
    for b in bullets:
        if BULLET_ICON.exists():
            slide.shapes.add_picture(str(BULLET_ICON), Inches(0.62), Inches(y),
                                     width=Inches(0.25), height=Inches(0.27))
        _add_textbox(slide, Inches(1.04), Inches(y - 0.02), Inches(5.2), Inches(0.45),
                     b, size=13.5, color=SLATE, font=FONT_BODY)
        y += 0.55

    # Right image
    img = image_path if image_path and image_path.exists() else SIDE_PANEL
    if img and img.exists():
        # Fit into right half
        slide.shapes.add_picture(str(img), Inches(6.55), Inches(1.40),
                                 width=Inches(6.2), height=Inches(5.4))
    return slide


def add_image_slide(prs, title: str, image_path: Path, caption: str = ""):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _title_with_accent(slide, title)

    if image_path.exists():
        # Leave room for caption
        top = Inches(1.35)
        max_h = Inches(5.0 if caption else 5.5)
        max_w = Inches(12.08)
        pic = slide.shapes.add_picture(str(image_path), Inches(0.62), top, width=max_w)
        # Scale down if too tall
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = max_h
            pic.left = int((Inches(13.333) - pic.width) / 2)

    if caption:
        _add_textbox(
            slide, Inches(0.62), Inches(6.55), Inches(12.08), Inches(0.45),
            caption, size=12, color=MUTED, font=FONT_BODY, align=PP_ALIGN.CENTER,
        )
    return slide


def add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]],
                    note: str = ""):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _title_with_accent(slide, title)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    left = Inches(0.62)
    top = Inches(1.55)
    width = Inches(12.08)
    row_h = Inches(0.58)
    height = row_h * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for j in range(n_cols):
        table.columns[j].width = int(width / n_cols)

    def _style_cell(cell, text, *, header=False, last=False):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        if header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY_BLUE
            _set_run_font(run, size=13, bold=True, color=WHITE, font=FONT_BODY)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if not last else RGBColor(0xEC, 0xFE, 0xFF)
            _set_run_font(run, size=13, bold=last, color=SLATE, font=FONT_BODY)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for j, h in enumerate(headers):
        _style_cell(table.cell(0, j), h, header=True)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            _style_cell(table.cell(i + 1, j), str(val), last=(i == len(rows) - 1))

    # Subtle row separators like template
    for i in range(1, n_rows):
        y = top + row_h * i
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, y, width, Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ROW_LINE
        line.line.fill.background()

    if note:
        _add_textbox(
            slide, Inches(0.62), Inches(1.55) + height + Inches(0.2),
            Inches(12.08), Inches(0.4),
            note, size=12, color=MUTED, font=FONT_BODY,
        )
    return slide


def add_highlight(prs, title: str, big: str, subtitle: str, bullets: list[str]):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _title_with_accent(slide, title)

    _add_textbox(
        slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(1.6),
        big, size=72, bold=True, color=CYAN, font=FONT_TITLE, align=PP_ALIGN.LEFT,
    )
    _add_textbox(
        slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(1.0),
        subtitle, size=16, color=SLATE, font=FONT_BODY,
    )
    y = 1.9
    for b in bullets:
        if BULLET_ICON.exists():
            slide.shapes.add_picture(str(BULLET_ICON), Inches(6.8), Inches(y),
                                     width=Inches(0.25), height=Inches(0.27))
        _add_textbox(slide, Inches(7.2), Inches(y - 0.02), Inches(5.4), Inches(0.55),
                     b, size=14, color=SLATE, font=FONT_BODY)
        y += 0.7
    return slide


def add_conclusion(prs, title: str, body: str):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONCLUSION)
    _add_textbox(
        slide, Inches(3.5), Inches(2.35), Inches(6.3), Inches(0.85),
        title, size=42, bold=True, color=WHITE, font=FONT_TITLE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, Inches(2.0), Inches(3.50), Inches(9.3), Inches(1.8),
        body, size=16, color=SLATE, font=FONT_BODY, align=PP_ALIGN.CENTER,
    )


def add_qa(prs):
    slide = _blank_slide(prs)
    _add_bg(slide, BG_CONTENT)
    _add_textbox(
        slide, Inches(2.5), Inches(2.35), Inches(8.3), Inches(1.1),
        "Hỏi & Đáp", size=54, bold=True, color=NAVY_BLUE, font=FONT_TITLE,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, Inches(2.5), Inches(3.70), Inches(8.3), Inches(0.55),
        "Xin cảm ơn Thầy/Cô và các bạn đã lắng nghe!",
        size=18, color=MUTED, font=FONT_BODY, align=PP_ALIGN.CENTER,
    )
    if ICON_QA.exists():
        slide.shapes.add_picture(str(ICON_QA), Inches(6.40), Inches(4.70),
                                 width=Inches(0.52), height=Inches(0.42))


def build():
    if not ASSETS.exists() or not BG_CONTENT.exists():
        raise FileNotFoundError(
            "Missing template-assets/. Re-extract media from the template first."
        )

    # Fresh deck with widescreen size matching the template (13.333 × 7.5 in).
    # Visual style comes from template background images + Merriweather/DM Sans.
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- 01 Cover ----
    add_cover(prs)

    # ---- Agenda as section-style list ----
    add_section(
        prs, "01", "Nội Dung Trình Bày",
        "Bối cảnh · Mục tiêu · Công nghệ · Kiến trúc · Hiện thực · Thử nghiệm · Kết quả · Kết luận",
    )

    add_bullets(
        prs, "Nội Dung Trình Bày",
        [
            "01 — Bối cảnh & vấn đề cold-start trong Serverless AI",
            "02 — Mục tiêu đề tài và quy trình 4 giai đoạn",
            "03 — Công nghệ & phương pháp Weight-Sharing (tmpfs)",
            "04 — Kiến trúc hệ thống trên Minikube + Knative",
            "05 — Hiện thực hóa: FastAPI, YOLO-World, YAML",
            "06 — Kịch bản thử nghiệm với k6 & Prometheus",
            "07 — Kết quả định lượng và phân tích",
            "08 — Kết luận, hạn chế và hướng phát triển",
        ],
    )

    # ---- 02 Background ----
    add_section(
        prs, "02", "Bối Cảnh & Vấn Đề",
        "AI + Serverless gặp rào cản cold-start khi nạp weights mô hình lớn.",
    )

    add_split_content(
        prs,
        "Bối Cảnh Đề Tài",
        "Serverless AI & Cold-start",
        "Deep Learning đòi hỏi hạ tầng linh hoạt, trong khi Serverless tối ưu chi phí qua scale-to-zero. Sự kết hợp này tạo ra nút thắt I/O khi khởi tạo container.",
        [
            "Mô hình weights lớn (YOLO-World ~90MB) nạp chậm từ đĩa.",
            "Mỗi lần scale-from-zero phát sinh cold-start đột biến.",
            "P99 latency suy giảm nghiêm trọng trải nghiệm thời gian thực.",
            "Cần tái thiết kế lớp lưu trữ dùng chung trên cùng node.",
        ],
        image_path=None,
    )

    add_bullets(
        prs, "Vấn Đề Cần Giải Quyết",
        [
            "Mỗi Pod mới phải đọc weights từ lưu trữ bền vững trên node → I/O đĩa là nút thắt.",
            "Câu hỏi nghiên cứu: Giảm chi phí nạp weights mà vẫn giữ kiến trúc Serverless?",
            "Hướng tiếp cận: Weight-Sharing qua phân vùng tmpfs trên RAM (hostPath).",
            "Đối chứng với Baseline đọc từ đĩa ảo trên cùng node Minikube.",
        ],
        intro="Đồ án tập trung vào tối ưu quá trình nạp trọng số mô hình — yếu tố then chốt của cold-start AI.",
    )

    # ---- 03 Objectives ----
    add_section(
        prs, "03", "Mục Tiêu Đề Tài",
        "Giảm đáng kể độ trễ cold-start khi nạp weights mô hình AI trên Serverless.",
    )

    add_three_cards(
        prs, "Mục Tiêu Cụ Thể",
        [
            ("Hạ tầng Serverless",
             "Thiết lập Minikube + Knative Serving với autoscaling và scale-to-zero tiêu chuẩn."),
            ("Weight-Sharing",
             "Dùng tmpfs 4GB trên RAM làm vùng lưu trữ dùng chung giữa các Pod qua hostPath."),
            ("Đo lường đối chứng",
             "k6 + Prometheus/Grafana: so sánh RAM vs Disk theo P99, Burst và Warm."),
        ],
    )

    add_image_slide(
        prs,
        "Quy Trình Thực Hiện (4 Giai Đoạn)",
        FIG / "00-quy-trinh-4-giai-doan.png",
        "Hạ tầng → Weight-Sharing → Load test → Tổng hợp & đánh giá",
    )

    # ---- 04 Technology ----
    add_section(
        prs, "04", "Công Nghệ & Phương Pháp",
        "Kubernetes, Knative, YOLO-World, FastAPI, k6 và Weight-Sharing qua tmpfs.",
    )

    add_two_cards(
        prs, "Stack Công Nghệ",
        "Hạ tầng & Điều phối",
        "Nền tảng Serverless trên cụm giả lập single-node.",
        [
            "Kubernetes (Minikube)",
            "Knative Serving (KPA, Activator)",
            "Kourier Ingress Gateway",
            "hostPath volume chia sẻ weights",
        ],
        "Ứng dụng & Giám sát",
        "Workload suy diễn và pipeline đo lường end-to-end.",
        [
            "YOLO-World (Ultralytics)",
            "FastAPI + Uvicorn + Prometheus",
            "Grafana dashboard đối chiếu",
            "k6 load test (cold / burst / cycle)",
        ],
    )

    add_two_cards(
        prs, "Phương Pháp Weight-Sharing",
        "Optimized — RAM tmpfs",
        "Một bản sao weights trên RAM; các Pod đọc chung qua hostPath.",
        [
            "tmpfs 4GB tại /mnt/shared-weights",
            "File yolov8l-world.pt trên RAM",
            "Giảm I/O đĩa khi cold-start",
            "P99 cold TB ≈ 17.35s",
        ],
        "Baseline — Disk",
        "Cùng file weights nhưng đọc từ đĩa ảo node — nhóm đối chứng.",
        [
            "Đĩa tại /mnt/disk-weights",
            "Mỗi cold-start đọc từ đĩa",
            "Nút thắt I/O rõ rệt",
            "P99 cold TB ≈ 75.00s",
        ],
    )

    add_image_slide(
        prs,
        "Weight-Sharing: tmpfs (RAM) vs Disk",
        FIG / "02-weight-sharing-tmpfs-vs-disk.png",
        "Các Pod trên cùng node chia sẻ một bản sao weights qua hostPath",
    )

    # ---- 05 Architecture ----
    add_section(
        prs, "05", "Kiến Trúc Hệ Thống",
        "Pipeline suy diễn AI Serverless trên Minikube với lớp lưu trữ dùng chung.",
    )

    add_image_slide(
        prs,
        "Kiến Trúc Tổng Thể",
        FIG / "01-kien-truc-pipeline-tong-the.png",
        "k6 → Kourier → Knative Serving → Pod FastAPI/YOLO · Prometheus/Grafana",
    )

    add_bullets(
        prs, "Các Thành Phần Chính",
        [
            "k6 — gửi POST /predict, đo end-to-end latency từ góc nhìn người dùng.",
            "Kourier — Ingress Gateway, định tuyến qua tên miền ảo *.sslip.io.",
            "Knative — Autoscaler (KPA), Activator (scale-from-zero), Queue-Proxy.",
            "Pod suy diễn — FastAPI + YOLO-World; nạp model tại module startup.",
            "Lớp lưu trữ — hostPath liên kết tmpfs (Optimized) hoặc disk (Baseline).",
            "Giám sát — Prometheus scrape /metrics; Grafana dashboard đối chiếu.",
        ],
    )

    # ---- 06 Implementation ----
    add_section(
        prs, "06", "Hiện Thực Hóa",
        "Hai dịch vụ Knative đối chứng, API FastAPI và tự động hóa bằng shell scripts.",
    )

    add_bullets(
        prs, "Cấu Hình Triển Khai",
        [
            "Hai dịch vụ độc lập: yolo-inference (optimized) và yolo-inference-baseline.",
            "Autoscaling: min-scale=0, max-scale=5, target concurrency=1.",
            "imagePullPolicy: Never — dùng image cục bộ, loại trừ biến số mạng.",
            "Prometheus annotations: scrape=true, port=8080, path=/metrics.",
            "Scripts: setup-weights.sh, deploy-all.sh — đảm bảo tái lập thí nghiệm.",
            "Model nạp top-level import → chi phí I/O xảy ra trước khi Pod báo Ready.",
        ],
    )

    add_bullets(
        prs, "API Cốt Lõi (FastAPI)",
        [
            "GET /health — trạng thái dịch vụ và deployment mode.",
            "POST /predict — nhận ảnh, chạy YOLO-World, trả JSON detections.",
            "GET /metrics — Prometheus: model load, latency histogram, request count.",
            "Metric chính: yolo_model_load_seconds, yolo_inference_request_duration_seconds.",
            "Docker image: python:3.10-slim + ultralytics + FastAPI + prometheus-client.",
        ],
    )

    # ---- 07 Experiments ----
    add_section(
        prs, "07", "Kịch Bản Thử Nghiệm",
        "Bốn kịch bản đo lường với chỉ số then chốt là P99 Latency.",
    )

    add_three_cards(
        prs, "Bốn Kịch Bản Đo Lường",
        [
            ("Cold-start",
             "3 VUs đồng thời, scale-to-zero trước mỗi lần chạy, lặp n=3 cho RAM và Disk."),
            ("Burst Traffic",
             "Ramp 15s → 15 VUs, giữ 60s, ramp-down 15s. Nghiệm thu: error < 10%, P99 < 60s."),
            ("Full Cycle",
             "Cold (3 VU) → Burst (15 VU) → Warm (5 VU, 1 phút) + giám sát model load."),
        ],
    )

    # ---- 08 Results ----
    add_section(
        prs, "08", "Kết Quả & Phân Tích",
        "Weight-Sharing qua RAM giảm 76.9% P99 cold-start so với Baseline đĩa.",
    )

    add_highlight(
        prs, "Kết Quả Nổi Bật",
        "−76.9%",
        "Giảm P99 cold-start trung bình\n(75.00s → 17.35s)",
        [
            "Optimized (RAM): P99 TB = 17.35s (σ = 1.68s)",
            "Baseline (Disk): P99 TB = 75.00s (σ = 12.78s)",
            "Burst: 0% lỗi HTTP · P99 = 46.37s (< 60s)",
            "Model load từ RAM: 1.38s cho file ~90MB",
            "Warm phase: P99 = 5.64s khi Pod sẵn sàng",
        ],
    )

    add_table_slide(
        prs,
        "P99 Cold-start (giây) — n=3",
        ["Lượt chạy", "Optimized (RAM)", "Baseline (Disk)"],
        [
            ["Lần 1", "15.45", "60.45"],
            ["Lần 2", "18.61", "80.12"],
            ["Lần 3", "18.00", "84.42"],
            ["Trung bình", "17.35", "75.00"],
        ],
        note="Độ lệch chuẩn: Optimized 1.68s · Baseline 12.78s",
    )

    add_image_slide(
        prs,
        "So Sánh P99 Cold-start",
        FIG / "03-cold-start-p99-so-sanh.png",
        "Giảm 76.9% P99 trung bình khi dùng Weight-Sharing qua RAM",
    )

    add_image_slide(
        prs,
        "Đối Chiếu P99 (Dashboard Style)",
        FIG / "07-grafana-style-p99-doi-chieu.png",
        "Optimized 17.35s vs Baseline 75.00s",
    )

    add_table_slide(
        prs,
        "Burst Traffic — Chỉ Số Hiệu Năng",
        ["Chỉ số", "Giá trị"],
        [
            ["Tổng requests", "61"],
            ["HTTP Error Rate", "0.00%"],
            ["Avg Latency", "19.40 s"],
            ["P99 Latency", "46.37 s"],
            ["Trạng thái", "Không tắc nghẽn"],
        ],
        note="Đáp ứng tiêu chí nghiệm thu: error < 10% và P99 < 60s",
    )

    add_image_slide(
        prs,
        "Phân Vị Độ Trễ Burst Traffic",
        FIG / "04-burst-latency-phan-vi.png",
        "P99 = 46.37s — dưới ngưỡng nghiệm thu 60s",
    )

    add_table_slide(
        prs,
        "Chu Kỳ Cold → Burst → Warm",
        ["Giai đoạn", "Requests", "Avg", "P99"],
        [
            ["1. Cold-start", "3", "19.75 s", "22.43 s"],
            ["2. Burst", "69", "17.71 s", "32.63 s"],
            ["3. Warm", "39", "3.64 s", "5.64 s"],
        ],
    )

    add_image_slide(
        prs,
        "Biến Động Độ Trễ Theo Giai Đoạn",
        FIG / "05-cold-burst-warm.png",
        "Warm P99 = 5.64s — khoảng cách Cold−Warm ước lượng chi phí khởi tạo Pod",
    )

    add_image_slide(
        prs,
        "Phân Tích Chi Phí Cold-start (Optimized)",
        FIG / "06-model-load-va-phan-bo-coldstart.png",
        "Nạp model 1.38s (~8%) — phần còn lại: scheduling, runtime Python, deserialize",
    )

    add_bullets(
        prs, "Thảo Luận",
        [
            "Weight-Sharing qua RAM loại bỏ hiệu quả nút thắt I/O đĩa khi nạp weights.",
            "Sau tối ưu I/O, cold-start còn ~17s chủ yếu do lập lịch Pod và runtime Python.",
            "Knative autoscaling xử lý burst tốt trong giới hạn max-scale=5.",
            "Pipeline k6 + Prometheus đo end-to-end và nội tại độc lập, có thể tái lập.",
            "Khoảng cách Cold (22.43s) vs Warm (5.64s) hỗ trợ thiết lập ngưỡng SLA.",
        ],
    )

    # ---- 09 Conclusion ----
    add_section(
        prs, "09", "Kết Luận & Hướng Phát Triển",
        "Tóm tắt đóng góp, hạn chế thực nghiệm và các hướng mở rộng tiếp theo.",
    )

    add_bullets(
        prs, "Kết Luận",
        [
            "Đã thiết kế và triển khai kiến trúc AI Serverless trên Kubernetes + Knative.",
            "Weight-Sharing qua tmpfs giảm 76.9% P99 cold-start (75.00s → 17.35s).",
            "Hệ thống ổn định dưới burst traffic: 0% lỗi HTTP, P99 = 46.37s < 60s.",
            "Xây dựng pipeline giám sát MLOps: k6 + Prometheus + Grafana.",
            "Đặt weights “nóng” trên RAM là chiến lược hiệu quả cho Serverless AI.",
        ],
    )

    add_two_cards(
        prs, "Hạn Chế & Hướng Phát Triển",
        "Hạn chế",
        "Phạm vi thực nghiệm còn giới hạn bởi môi trường giả lập.",
        [
            "Minikube single-node — chưa đa node",
            "tmpfs 4GB cấu hình tĩnh",
            "Cold-start còn ~17s sau tối ưu I/O",
            "Mẫu thống kê n=3 — mang tính minh họa",
        ],
        "Hướng phát triển",
        "Mở rộng sang môi trường phân tán và tối ưu runtime sâu hơn.",
        [
            "Cache phân tán: Alluxio multi-node",
            "Nén model: Quantization / Pruning",
            "Pre-warming theo dự báo lưu lượng",
            "Engine: Triton / ONNX Runtime",
        ],
    )

    add_conclusion(
        prs,
        "Kết Luận",
        "Đồ án đã giảm đáng kể nút thắt I/O đĩa khi nạp trọng số AI trên hạ tầng Serverless. "
        "Weight-Sharing qua RAM là chiến lược hiệu quả; chi phí lập lịch và runtime vẫn là "
        "giới hạn tiếp theo cần tối ưu.",
    )

    add_qa(prs)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
