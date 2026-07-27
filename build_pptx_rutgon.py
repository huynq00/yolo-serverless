#!/usr/bin/env python3
"""Build a streamlined ~17-slide presentation from the YOLO Serverless report."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
MEDIA = ROOT / "_pptx_media_tmp"
FIGS = ROOT / "report-figures" / "PNG-copy-vao-bao-cao"
OUT = ROOT / "Thuyet_Trinh_YOLO_Serverless_RutGon.pptx"

# Theme (match original dark-blue deck)
BG = RGBColor(0x0B, 0x1F, 0x3F)
BG_CARD = RGBColor(0x12, 0x2E, 0x56)
ACCENT = RGBColor(0x3D, 0xB8, 0xE8)
ACCENT2 = RGBColor(0x5E, 0xE0, 0xA8)
WARN = RGBColor(0xF0, 0xA5, 0x4A)
TEXT = RGBColor(0xF2, 0xF6, 0xFC)
MUTED = RGBColor(0xA8, 0xB8, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE8, 0x6A, 0x6A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, size=18, bold=False, color=TEXT, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, page, total=17):
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.05), Inches(10), Inches(0.35)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "YOLO Serverless · Weight-Sharing · NT2204.CH201"
    set_run(run, 11, color=MUTED)
    num = slide.shapes.add_textbox(
        Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.35)
    )
    tf2 = num.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = f"{page}/{total}"
    set_run(run2, 11, color=MUTED)


def add_title(slide, text, top=0.28, size=28):
    box = slide.shapes.add_textbox(
        Inches(0.55), Inches(top), Inches(12.2), Inches(0.7)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size, bold=True, color=WHITE)
    return box


def add_subtitle(slide, text, top=0.85, size=15):
    box = slide.shapes.add_textbox(
        Inches(0.55), Inches(top), Inches(12.2), Inches(0.4)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size, color=ACCENT)
    return box


def add_bullets(slide, items, left=0.55, top=1.4, width=12.0, height=5.2, size=18):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size, color=TEXT)
    return box


def add_card(slide, left, top, width, height, fill=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0x1E, 0x45, 0x78)
    shape.line.width = Pt(1)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_textbox(slide, left, top, width, height, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold=bold, color=color)
    return box


def add_picture_fit(slide, path, left, top, max_w, max_h):
    path = Path(path)
    if not path.exists():
        add_textbox(slide, left, top, max_w, 0.4, f"[Thiếu hình: {path.name}]", 12, color=WARN)
        return None
    from PIL import Image

    with Image.open(path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px
    box_aspect = max_w / max_h
    if aspect >= box_aspect:
        w = max_w
        h = max_w / aspect
    else:
        h = max_h
        w = max_h * aspect
    # center in box
    x = left + (max_w - w) / 2
    y = top + (max_h - h) / 2
    return slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), Inches(w), Inches(h)
    )


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 17

    # ---- 1. Title ----
    s = new_slide(prs)
    title_img = MEDIA / "image1.png"
    if title_img.exists():
        s.shapes.add_picture(str(title_img), 0, 0, SLIDE_W, SLIDE_H)
    else:
        add_bg(s)
    add_accent_bar(s)
    add_textbox(s, 0.7, 1.6, 11.5, 0.45, "BÁO CÁO ĐỒ ÁN", size=14, bold=True, color=ACCENT)
    add_textbox(
        s, 0.7, 2.2, 11.8, 1.4,
        "Tối ưu hóa Cold-Start AI Serverless\nbằng Weight-Sharing",
        size=32, bold=True, color=WHITE,
    )
    add_textbox(
        s, 0.7, 4.0, 11.5, 0.45,
        "Hệ tính toán phân bố nâng cao (NT2204.CH201)",
        size=16, color=MUTED,
    )
    add_textbox(
        s, 0.7, 5.2, 11.5, 0.4,
        "Sinh viên thực hiện: Ngô Quang Huy",
        size=15, color=TEXT,
    )
    add_textbox(
        s, 0.7, 5.65, 11.5, 0.4,
        "Giảng viên hướng dẫn: TS. Huỳnh Văn Đặng",
        size=15, color=TEXT,
    )
    add_textbox(s, 11.5, 6.9, 1.4, 0.35, "1/17", size=12, color=MUTED, align=PP_ALIGN.RIGHT)

    # ---- 2. Agenda ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 2, total)
    add_title(s, "Nội dung trình bày")
    add_subtitle(s, "Một mạch: vấn đề → giải pháp → đo lường → kết quả")
    agenda = [
        ("01", "Vấn đề & câu hỏi nghiên cứu"),
        ("02", "Mục tiêu và ý tưởng Weight-Sharing"),
        ("03", "Kiến trúc & hiện thực rút gọn"),
        ("04", "Kịch bản thử nghiệm"),
        ("05", "Kết quả định lượng & thảo luận"),
        ("06", "Kết luận, hạn chế, hướng phát triển"),
    ]
    for i, (num, label) in enumerate(agenda):
        row = i // 2
        col = i % 2
        left = 0.55 + col * 6.3
        top = 1.5 + row * 1.5
        add_card(s, left, top, 5.9, 1.25)
        add_textbox(s, left + 0.25, top + 0.25, 1.0, 0.5, num, size=22, bold=True, color=ACCENT)
        add_textbox(s, left + 1.3, top + 0.35, 4.3, 0.55, label, size=16, color=TEXT)

    # ---- 3. Problem + research question ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 3, total)
    add_title(s, "Bối cảnh & câu hỏi nghiên cứu")
    add_subtitle(s, "Serverless AI gặp nút thắt I/O khi nạp weights lúc scale-from-zero")
    add_card(s, 0.55, 1.45, 7.6, 4.8)
    add_bullets(
        s,
        [
            "Scale-to-zero tiết kiệm chi phí, nhưng mỗi Pod mới phải cold-start.",
            "YOLO-World ~90MB: nạp weights từ đĩa là nút thắt chính.",
            "P99 latency tăng đột biến → khó đáp ứng thời gian thực.",
        ],
        left=0.8, top=1.7, width=7.1, height=2.6, size=17,
    )
    add_textbox(s, 0.8, 4.4, 7.1, 0.35, "CÂU HỎI NGHIÊN CỨU", size=12, bold=True, color=ACCENT)
    add_textbox(
        s, 0.8, 4.85, 7.1, 1.1,
        "Làm sao giảm chi phí nạp weights mà vẫn giữ kiến trúc Serverless (autoscaling + scale-to-zero)?",
        size=18, bold=True, color=WHITE,
    )
    add_card(s, 8.4, 1.45, 4.4, 4.8)
    add_textbox(s, 8.65, 1.7, 4.0, 0.4, "Hướng tiếp cận", size=14, bold=True, color=ACCENT)
    add_bullets(
        s,
        [
            "Weight-Sharing qua tmpfs (RAM)",
            "hostPath dùng chung trên node",
            "Đối chứng Baseline đọc từ đĩa",
        ],
        left=8.65, top=2.3, width=3.9, height=3.5, size=16,
    )

    # ---- 4. Goals ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 4, total)
    add_title(s, "Mục tiêu cụ thể")
    add_subtitle(s, "Ba mục tiêu đo được trong phạm vi đồ án")
    goals = [
        ("Hạ tầng Serverless", "Minikube + Knative Serving\nautoscaling & scale-to-zero"),
        ("Weight-Sharing", "tmpfs 4GB trên RAM\nchia sẻ weights qua hostPath"),
        ("Đo lường đối chứng", "k6 + Prometheus/Grafana\nP99: RAM vs Disk"),
    ]
    for i, (title, body) in enumerate(goals):
        left = 0.55 + i * 4.15
        add_card(s, left, 1.6, 3.95, 4.5)
        add_textbox(s, left + 0.25, 1.9, 3.45, 0.4, f"0{i+1}", size=20, bold=True, color=ACCENT)
        add_textbox(s, left + 0.25, 2.5, 3.45, 0.8, title, size=18, bold=True, color=WHITE)
        add_textbox(s, left + 0.25, 3.5, 3.45, 2.0, body, size=16, color=MUTED)

    # ---- 5. Weight-Sharing idea (before/after) ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 5, total)
    add_title(s, "Ý tưởng: Weight-Sharing (tmpfs vs Disk)", size=24)
    fig = FIGS / "Hinh-3.2-weight-sharing-tmpfs-vs-disk.png"
    if not fig.exists():
        fig = MEDIA / "image8.png"
    add_picture_fit(s, fig, 0.4, 1.05, 12.5, 5.8)

    # ---- 6. Stack short ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 6, total)
    add_title(s, "Stack công nghệ (rút gọn)")
    add_subtitle(s, "Chỉ các thành phần cần để hiểu pipeline end-to-end")
    left_items = [
        ("Hạ tầng", "Kubernetes (Minikube), Knative Serving, Kourier"),
        ("Ứng dụng", "FastAPI + YOLO-World (Ultralytics)"),
        ("Lưu trữ", "hostPath → tmpfs (Optimized) / disk (Baseline)"),
        ("Đo lường", "k6 · Prometheus · Grafana"),
    ]
    for i, (h, b) in enumerate(left_items):
        top = 1.45 + i * 1.2
        add_card(s, 0.55, top, 12.2, 1.05)
        add_textbox(s, 0.85, top + 0.28, 2.4, 0.5, h, size=16, bold=True, color=ACCENT)
        add_textbox(s, 3.4, top + 0.28, 9.0, 0.5, b, size=16, color=TEXT)

    # ---- 7. Architecture ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 7, total)
    add_title(s, "Kiến trúc tổng thể", size=24)
    fig = FIGS / "Hinh-3.1-pipeline-kien-truc-tong-the.png"
    if not fig.exists():
        fig = MEDIA / "image9.png"
    add_picture_fit(s, fig, 0.35, 1.0, 12.6, 5.85)

    # ---- 8. Implementation short ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 8, total)
    add_title(s, "Hiện thực rút gọn")
    add_subtitle(s, "Hai dịch vụ đối chứng — cùng image, khác lớp lưu trữ weights")
    add_card(s, 0.55, 1.5, 6.0, 4.7)
    add_textbox(s, 0.85, 1.75, 5.4, 0.4, "Hai Knative Service", size=16, bold=True, color=ACCENT)
    add_bullets(
        s,
        [
            "yolo-inference → tmpfs (Optimized)",
            "yolo-inference-baseline → disk",
            "min-scale=0, max-scale=5, concurrency=1",
            "Model nạp lúc startup → I/O vào cold-start",
        ],
        left=0.85, top=2.35, width=5.4, height=3.5, size=16,
    )
    add_card(s, 6.8, 1.5, 5.95, 4.7)
    add_textbox(s, 7.1, 1.75, 5.4, 0.4, "API đo lường", size=16, bold=True, color=ACCENT)
    add_bullets(
        s,
        [
            "POST /predict — suy diễn YOLO",
            "GET /metrics — Prometheus",
            "yolo_model_load_seconds",
            "yolo_inference_request_duration_seconds",
        ],
        left=7.1, top=2.35, width=5.4, height=3.5, size=16,
    )

    # ---- 9. Experiment scenarios ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 9, total)
    add_title(s, "Kịch bản thử nghiệm")
    add_subtitle(s, "Chỉ số then chốt: P99 latency (góc nhìn người dùng qua k6)")
    scenarios = [
        ("Cold-start", "3 VU · scale-to-zero trước mỗi lần\nlặp n=3 cho RAM và Disk"),
        ("Burst", "Ramp → 15 VU · 60s\nnghiệm thu: error <10%, P99 <60s"),
        ("Full cycle", "Cold → Burst → Warm\nước lượng chi phí khởi tạo Pod"),
    ]
    for i, (t, b) in enumerate(scenarios):
        left = 0.55 + i * 4.15
        add_card(s, left, 1.6, 3.95, 4.4)
        add_textbox(s, left + 0.3, 2.0, 3.4, 0.5, t, size=20, bold=True, color=ACCENT2 if i == 0 else WHITE)
        add_textbox(s, left + 0.3, 2.9, 3.4, 2.5, b, size=16, color=MUTED)

    # ---- 10. Headline result ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 10, total)
    add_title(s, "Kết quả nổi bật")
    add_subtitle(s, "Weight-Sharing qua RAM loại bỏ nút thắt I/O đĩa khi cold-start")
    add_card(s, 0.55, 1.55, 5.5, 4.7)
    add_textbox(s, 0.9, 2.0, 4.8, 0.5, "Giảm P99 cold-start", size=16, color=ACCENT)
    add_textbox(s, 0.9, 2.6, 4.8, 1.0, "−76.9%", size=60, bold=True, color=ACCENT2)
    add_textbox(s, 0.9, 3.9, 4.8, 0.5, "75.00s  →  17.35s", size=24, bold=True, color=WHITE)
    add_textbox(s, 0.9, 4.7, 4.8, 0.8, "Optimized σ=1.68s  ·  Baseline σ=12.78s", size=14, color=MUTED)
    add_card(s, 6.3, 1.55, 6.45, 4.7)
    add_textbox(s, 6.65, 1.85, 5.9, 0.4, "Các chỉ số kèm theo", size=15, bold=True, color=ACCENT)
    add_bullets(
        s,
        [
            "Model load từ RAM: 1.38s (~90MB)",
            "Burst: 0% lỗi HTTP · P99 = 46.37s (< 60s)",
            "Warm phase: P99 = 5.64s khi Pod sẵn sàng",
            "Khoảng cách Cold−Warm ≈ chi phí khởi tạo",
        ],
        left=6.65, top=2.5, width=5.8, height=3.4, size=16,
    )

    # ---- 11. P99 chart ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 11, total)
    add_title(s, "So sánh P99 cold-start (n=3)")
    add_subtitle(s, "Optimized ổn định hơn Baseline (σ thấp hơn rõ rệt)")
    fig = FIGS / "Hinh-4.1-cold-start-p99-so-sanh.png"
    if not fig.exists():
        fig = MEDIA / "image10.png"
    add_picture_fit(s, fig, 0.5, 1.3, 8.5, 5.3)
    # side table cards
    add_card(s, 9.2, 1.5, 3.6, 2.1)
    add_textbox(s, 9.4, 1.7, 3.2, 0.35, "Optimized (RAM)", size=13, bold=True, color=ACCENT2)
    add_textbox(s, 9.4, 2.2, 3.2, 0.8, "17.35s", size=32, bold=True, color=WHITE)
    add_textbox(s, 9.4, 3.05, 3.2, 0.35, "15.45 · 18.61 · 18.00", size=12, color=MUTED)
    add_card(s, 9.2, 3.9, 3.6, 2.1)
    add_textbox(s, 9.4, 4.1, 3.2, 0.35, "Baseline (Disk)", size=13, bold=True, color=WARN)
    add_textbox(s, 9.4, 4.6, 3.2, 0.8, "75.00s", size=32, bold=True, color=WHITE)
    add_textbox(s, 9.4, 5.45, 3.2, 0.35, "60.45 · 80.12 · 84.42", size=12, color=MUTED)

    # ---- 12. Burst + full cycle ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 12, total)
    add_title(s, "Burst traffic & chu kỳ Cold → Burst → Warm")
    add_subtitle(s, "Hệ thống ổn định dưới tải; Warm phản ánh latency khi Pod sẵn sàng")
    # left burst metrics
    add_card(s, 0.55, 1.45, 5.9, 5.0)
    add_textbox(s, 0.85, 1.7, 5.3, 0.4, "Burst — nghiệm thu đạt", size=16, bold=True, color=ACCENT)
    metrics = [
        ("Tổng requests", "61"),
        ("HTTP error", "0.00%"),
        ("Avg latency", "19.40 s"),
        ("P99 latency", "46.37 s  (< 60s)"),
    ]
    for i, (k, v) in enumerate(metrics):
        top = 2.3 + i * 0.9
        add_textbox(s, 0.95, top, 2.8, 0.4, k, size=15, color=MUTED)
        add_textbox(s, 3.6, top, 2.5, 0.4, v, size=16, bold=True, color=WHITE)
    # right cycle table
    add_card(s, 6.7, 1.45, 6.05, 5.0)
    add_textbox(s, 7.0, 1.7, 5.5, 0.4, "Full cycle (Optimized)", size=16, bold=True, color=ACCENT)
    headers = [("Giai đoạn", 7.0), ("Req", 9.3), ("Avg", 10.3), ("P99", 11.4)]
    for h, x in headers:
        add_textbox(s, x, 2.35, 1.3, 0.35, h, size=13, bold=True, color=MUTED)
    rows = [
        ("1. Cold-start", "3", "19.75s", "22.43s"),
        ("2. Burst", "69", "17.71s", "32.63s"),
        ("3. Warm", "39", "3.64s", "5.64s"),
    ]
    for i, row in enumerate(rows):
        top = 2.9 + i * 0.85
        colors = [TEXT, TEXT, TEXT, ACCENT2 if i == 2 else TEXT]
        xs = [7.0, 9.3, 10.3, 11.4]
        for j, (val, x) in enumerate(zip(row, xs)):
            add_textbox(s, x, top, 1.4, 0.4, val, size=15, bold=(j == 0 or j == 3), color=colors[j])

    # ---- 13. Cost breakdown ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 13, total)
    add_title(s, "Sau khi tối ưu I/O — còn gì trong cold-start?", size=24)
    fig = FIGS / "Hinh-4.4-model-load-va-phan-bo-coldstart.png"
    if not fig.exists():
        fig = MEDIA / "image14.png"
    add_picture_fit(s, fig, 0.35, 1.0, 12.6, 5.85)

    # ---- 14. Discussion ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 14, total)
    add_title(s, "Thảo luận")
    add_subtitle(s, "Ý nghĩa kết quả đối với thiết kế Serverless AI")
    points = [
        ("I/O đã hết là nút thắt chính", "tmpfs loại bỏ hiệu quả chi phí đọc weights từ đĩa."),
        ("~17s còn lại là giới hạn mới", "Chủ yếu scheduling Pod + khởi động runtime Python."),
        ("Autoscaling xử lý burst tốt", "Trong giới hạn max-scale=5, P99 vẫn dưới ngưỡng nghiệm thu."),
        ("Cold vs Warm hỗ trợ SLA", "Khoảng cách latency giúp đặt ngưỡng / chiến lược pre-warm."),
    ]
    for i, (h, b) in enumerate(points):
        row, col = divmod(i, 2)
        left = 0.55 + col * 6.3
        top = 1.5 + row * 2.4
        add_card(s, left, top, 6.05, 2.15)
        add_textbox(s, left + 0.3, top + 0.35, 5.45, 0.5, h, size=16, bold=True, color=ACCENT)
        add_textbox(s, left + 0.3, top + 1.0, 5.45, 0.8, b, size=15, color=TEXT)

    # ---- 15. Limitations & future ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 15, total)
    add_title(s, "Hạn chế & hướng phát triển")
    add_card(s, 0.55, 1.45, 6.0, 5.0)
    add_textbox(s, 0.85, 1.7, 5.4, 0.4, "Hạn chế", size=18, bold=True, color=WARN)
    add_bullets(
        s,
        [
            "Minikube single-node — chưa đa node",
            "tmpfs 4GB cấu hình tĩnh",
            "Cold-start còn ~17s sau tối ưu I/O",
            "Mẫu n=3 mang tính minh họa",
        ],
        left=0.85, top=2.4, width=5.4, height=3.6, size=16,
    )
    add_card(s, 6.8, 1.45, 5.95, 5.0)
    add_textbox(s, 7.1, 1.7, 5.4, 0.4, "Hướng phát triển", size=18, bold=True, color=ACCENT2)
    add_bullets(
        s,
        [
            "Cache phân tán: Alluxio multi-node",
            "Nén model: Quantization / Pruning",
            "Pre-warming theo dự báo lưu lượng",
            "Engine: Triton / ONNX Runtime",
        ],
        left=7.1, top=2.4, width=5.4, height=3.6, size=16,
    )

    # ---- 16. Conclusion ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s); add_footer(s, 16, total)
    add_title(s, "Kết luận — take-home")
    add_subtitle(s, "Ba điểm cần nhớ khi rời phòng")
    takeaways = [
        ("1", "Weight-Sharing qua RAM giảm 76.9% P99 cold-start\n(75.00s → 17.35s) trên cùng hạ tầng Serverless."),
        ("2", "Sau tối ưu I/O, nút thắt chuyển sang scheduling\nvà runtime Python — hướng tối ưu tiếp theo."),
        ("3", "Pipeline đo k6 + Prometheus tái lập được,\nhỗ trợ đánh giá SLA Cold / Warm rõ ràng."),
    ]
    for i, (num, text) in enumerate(takeaways):
        top = 1.5 + i * 1.6
        add_card(s, 0.55, top, 12.2, 1.4)
        add_textbox(s, 0.85, top + 0.35, 0.8, 0.6, num, size=28, bold=True, color=ACCENT)
        add_textbox(s, 1.8, top + 0.3, 10.5, 0.9, text, size=17, color=TEXT)

    # ---- 17. Q&A ----
    s = new_slide(prs)
    add_bg(s); add_accent_bar(s)
    add_textbox(
        s, 0.8, 2.4, 11.5, 1.0,
        "Hỏi & Đáp",
        size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, 0.8, 3.6, 11.5, 0.6,
        "Xin cảm ơn Thầy/Cô và các bạn đã lắng nghe!",
        size=20, color=ACCENT, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, 0.8, 4.5, 11.5, 0.5,
        "Ngô Quang Huy  ·  NT2204.CH201  ·  YOLO Serverless Weight-Sharing",
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )
    add_footer(s, 17, total)

    prs.save(OUT)
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes), slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
